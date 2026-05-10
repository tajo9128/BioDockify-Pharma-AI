from __future__ import annotations

import io
import json
import re
import zipfile
from typing import Any
from xml.sax.saxutils import escape


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def pptx_from_text(title: str, content: str) -> bytes:
    return pptx_from_slides(slides_from_text(title, content))


def slides_from_text(title: str, content: str) -> list[dict[str, Any]]:
    normalized = normalize_slides(content)
    document_title = str(title or "Presentation").strip() or "Presentation"
    if not normalized:
        return [{"title": document_title, "bullets": []}]

    text = str(content or "")
    if len(normalized) == 1 and "---" not in text:
        slide = normalized[0]
        slide_title = str(slide.get("title") or "").strip()
        if slide_title and slide_title.casefold() != document_title.casefold():
            return [{"title": document_title, "bullets": [slide_title, *slide.get("bullets", [])]}]
    return normalized


def normalize_slides(value: Any) -> list[dict[str, Any]]:
    if value is None:
        return []
    if isinstance(value, str):
        stripped = value.strip()
        if not stripped:
            return []
        if stripped.startswith("[") or stripped.startswith("{"):
            return normalize_slides(json.loads(stripped))
        chunks = re.split(r"(?m)^\s*---+\s*$", stripped)
        result = []
        for chunk in chunks:
            lines = [_clean_slide_line(line) for line in chunk.splitlines() if line.strip()]
            lines = [line for line in lines if line]
            if lines:
                result.append({"title": lines[0], "bullets": lines[1:]})
        return result
    if isinstance(value, dict):
        return [_slide_from_mapping(value)]
    if isinstance(value, list):
        result = []
        for item in value:
            if isinstance(item, dict):
                result.append(_slide_from_mapping(item))
            elif isinstance(item, str):
                result.extend(normalize_slides(item))
            elif isinstance(item, (list, tuple)):
                lines = [_clean_slide_line(part) for part in item if str(part).strip()]
                if lines:
                    result.append({"title": lines[0], "bullets": lines[1:]})
            else:
                result.append({"title": str(item), "bullets": []})
        return result
    return [{"title": str(value), "bullets": []}]


def pptx_from_slides(slides: list[dict[str, Any]]) -> bytes:
    normalized = normalize_slides(slides)
    if not normalized:
        normalized = [{"title": "Presentation", "bullets": []}]
    try:
        return _pptx_from_slides_with_python_pptx(normalized)
    except Exception:
        return _pptx_from_slides_ooxml(normalized)


def _slide_from_mapping(value: dict[str, Any]) -> dict[str, Any]:
    title = _clean_slide_line(value.get("title") or value.get("heading") or "Slide")
    bullets = value.get("bullets")
    if bullets is None:
        body = value.get("body") or value.get("content") or ""
        bullets = [_clean_slide_line(line) for line in str(body).splitlines() if line.strip()]
    elif isinstance(bullets, str):
        bullets = [_clean_slide_line(line) for line in bullets.splitlines() if line.strip()]
    else:
        bullets = [_clean_slide_line(item) for item in bullets]
    return {"title": title or "Slide", "bullets": [bullet for bullet in bullets if bullet]}


def _clean_slide_line(value: Any) -> str:
    line = str(value or "").strip()
    line = re.sub(r"^\s{0,3}#{1,6}\s+", "", line)
    line = re.sub(r"^\s*(?:[-*•]|\d+[.)])\s+", "", line)
    return line.strip()


def _pptx_from_slides_with_python_pptx(slides: list[dict[str, Any]]) -> bytes:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches  # type: ignore

    presentation = Presentation()
    for slide_spec in slides:
        layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        title = str(slide_spec.get("title") or "Slide")
        bullets = [str(item) for item in slide_spec.get("bullets") or []]

        if slide.shapes.title:
            slide.shapes.title.text = title
        else:
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.8))
            title_box.text_frame.text = title

        body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_shape is None:
            body_shape = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(8.35), Inches(4.55))

        text_frame = body_shape.text_frame
        text_frame.clear()
        if not bullets:
            text_frame.text = ""
            continue

        for index, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _pptx_from_slides_ooxml(slides: list[dict[str, Any]]) -> bytes:
    files: dict[str, str | bytes] = {
        "[Content_Types].xml": _pptx_content_types(len(slides)),
        "_rels/.rels": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<Relationships xmlns="{REL_NS}">'
            '<Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
            'Target="ppt/presentation.xml"/>'
            "</Relationships>"
        ),
        "ppt/_rels/presentation.xml.rels": _pptx_presentation_rels(len(slides)),
        "ppt/presentation.xml": _pptx_presentation_xml(len(slides)),
    }
    for index, slide in enumerate(slides, start=1):
        files[f"ppt/slides/slide{index}.xml"] = _pptx_slide_xml(slide)
    return _zip_map(files)


def _pptx_content_types(count: int) -> str:
    overrides = [
        '<Override PartName="/ppt/presentation.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
    ]
    for index in range(1, count + 1):
        overrides.append(
            f'<Override PartName="/ppt/slides/slide{index}.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        )
    return (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<Types xmlns="{CT_NS}">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        + "".join(overrides)
        + "</Types>"
    )


def _pptx_presentation_rels(count: int) -> str:
    rels = []
    for index in range(1, count + 1):
        rels.append(
            f'<Relationship Id="rId{index}" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
            f'Target="slides/slide{index}.xml"/>'
        )
    return '<?xml version="1.0" encoding="UTF-8"?>' + f'<Relationships xmlns="{REL_NS}">' + "".join(rels) + "</Relationships>"


def _pptx_presentation_xml(count: int) -> str:
    slide_ids = "".join(f'<p:sldId id="{255 + index}" r:id="rId{index}"/>' for index in range(1, count + 1))
    return (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<p:presentation xmlns:p="{P_NS}" xmlns:r="{R_NS}">'
        f"<p:sldIdLst>{slide_ids}</p:sldIdLst>"
        '<p:sldSz cx="9144000" cy="5143500"/>'
        "</p:presentation>"
    )


def _pptx_slide_xml(slide: dict[str, Any]) -> str:
    title = str(slide.get("title") or "Slide")
    bullets = [str(item) for item in slide.get("bullets") or []]
    paragraphs = [title, *bullets]
    text = "".join(f"<a:p><a:r><a:t>{escape(item)}</a:t></a:r></a:p>" for item in paragraphs)
    return (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<p:sld xmlns:a="{A_NS}" xmlns:p="{P_NS}">'
        "<p:cSld><p:spTree>"
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        "<p:grpSpPr/>"
        '<p:sp><p:nvSpPr><p:cNvPr id="2" name="Content"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f"<p:txBody><a:bodyPr/><a:lstStyle/>{text}</p:txBody>"
        "</p:sp>"
        "</p:spTree></p:cSld>"
        "</p:sld>"
    )


def _zip_map(files_map: dict[str, str | bytes]) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, value in files_map.items():
            archive.writestr(name, value.encode("utf-8") if isinstance(value, str) else value)
    return buffer.getvalue()
